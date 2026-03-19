"""Document format conversion and validation utilities."""
from __future__ import annotations

import shutil
import subprocess
from pathlib import Path

SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".pptx", ".xlsx"}

# Known LibreOffice install paths on Windows
_SOFFICE_CANDIDATES = [
    "soffice",
    r"C:\Program Files\LibreOffice\program\soffice.exe",
    r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
]


def _find_soffice() -> str:
    """Find the soffice executable."""
    for candidate in _SOFFICE_CANDIDATES:
        if shutil.which(candidate) or Path(candidate).is_file():
            return candidate
    raise RuntimeError(
        "LibreOffice not found. Install it or add soffice to PATH."
    )

# Excel limits
MAX_SHEETS = 50
MAX_ROWS = 100000
MAX_COLS = 100


def is_supported_format(filepath: str | Path) -> bool:
    """Check whether the file extension is in the supported set."""
    return Path(filepath).suffix.lower() in SUPPORTED_EXTENSIONS


def check_excel_limits(filepath: str | Path) -> tuple[bool, str]:
    """Validate an Excel file against sheet/row/column limits.

    Returns (True, "") if within limits, or (False, reason) otherwise.
    """
    import openpyxl

    wb = openpyxl.load_workbook(filepath, read_only=True, data_only=True)
    try:
        sheet_names = wb.sheetnames
        if len(sheet_names) > MAX_SHEETS:
            return False, f"Too many sheets: {len(sheet_names)} (max {MAX_SHEETS})"

        for name in sheet_names:
            ws = wb[name]
            if ws.max_row is not None and ws.max_row > MAX_ROWS:
                return False, (
                    f"Sheet '{name}' has too many rows: {ws.max_row} (max {MAX_ROWS})"
                )
            if ws.max_column is not None and ws.max_column > MAX_COLS:
                return False, (
                    f"Sheet '{name}' has too many columns: {ws.max_column} (max {MAX_COLS})"
                )

        return True, ""
    finally:
        wb.close()


def convert_to_pdf(input_path: str | Path, output_dir: str | Path) -> Path:
    """Convert a document to PDF using LibreOffice headless mode.

    Returns the path to the generated PDF file.
    Raises RuntimeError if the conversion fails.
    """
    input_path = Path(input_path)
    output_dir = Path(output_dir)

    soffice = _find_soffice()
    result = subprocess.run(
        [
            soffice,
            "--headless",
            "--convert-to",
            "pdf",
            "--outdir",
            str(output_dir),
            str(input_path),
        ],
        capture_output=True,
        text=True,
    )

    if result.returncode != 0:
        raise RuntimeError(
            f"LibreOffice conversion failed (rc={result.returncode}): {result.stderr}"
        )

    pdf_path = output_dir / f"{input_path.stem}.pdf"
    if not pdf_path.exists():
        raise RuntimeError(
            f"Conversion appeared to succeed but PDF not found at {pdf_path}"
        )

    return pdf_path


def get_page_count(pdf_path: str | Path) -> int:
    """Return the number of pages in a PDF file."""
    from pypdf import PdfReader

    reader = PdfReader(pdf_path)
    return len(reader.pages)


def extract_pages_native(filepath: str | Path) -> list[str]:
    """Extract text per page/slide/sheet directly from the source file.

    Returns a list of strings, one per logical page.
    Falls back to PDF extraction if native extraction is not available.
    """
    filepath = Path(filepath)
    fmt = filepath.suffix.lower()

    if fmt in (".pptx",):
        return _extract_pptx(filepath)
    if fmt in (".docx",):
        return _extract_docx(filepath)
    if fmt in (".xlsx",):
        return _extract_xlsx(filepath)
    if fmt == ".pdf":
        return _extract_pdf(filepath)
    return []


def _extract_pptx(filepath: Path) -> list[str]:
    from pptx import Presentation

    prs = Presentation(filepath)
    pages = []
    for slide in prs.slides:
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if t:
                    texts.append(t)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if cells:
                        texts.append(" | ".join(cells))
        pages.append("\n".join(texts))
    return pages


def _extract_docx(filepath: Path) -> list[str]:
    """Extract text from docx. Treats the whole document as one page."""
    from docx import Document

    doc = Document(filepath)
    paragraphs = []
    for para in doc.paragraphs:
        t = para.text.strip()
        if t:
            paragraphs.append(t)
    # Also extract tables
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                paragraphs.append(" | ".join(cells))
    # Split into chunks of ~2000 chars to create logical pages
    full = "\n".join(paragraphs)
    if not full:
        return []
    return _chunk_text(full, 2000)


def _extract_xlsx(filepath: Path) -> list[str]:
    import openpyxl

    wb = openpyxl.load_workbook(filepath, read_only=True, data_only=True)
    sheet_names = wb.sheetnames
    wb.close()

    # Extract AutoShape text per sheet from drawing XML
    shape_texts = _extract_xlsx_shapes(filepath, len(sheet_names))

    wb = openpyxl.load_workbook(filepath, read_only=True, data_only=True)
    pages = []
    try:
        for idx, name in enumerate(wb.sheetnames):
            ws = wb[name]
            lines = [f"[Sheet: {name}]"]
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                line = " | ".join(cells).strip()
                if line and line != "| " * (len(cells) - 1):
                    lines.append(line)
            # Append AutoShape text if any
            if idx in shape_texts and shape_texts[idx]:
                lines.append("[AutoShape]")
                lines.append(shape_texts[idx])
            sheet_text = "\n".join(lines)
            # Large sheets: chunk into ~2000 char pages for better search
            if len(sheet_text) > 3000:
                chunks = _chunk_text(sheet_text, 2000)
                pages.extend(chunks)
            else:
                pages.append(sheet_text)
    finally:
        wb.close()
    return pages


def _extract_xlsx_shapes(filepath: Path, sheet_count: int) -> dict[int, str]:
    """Extract text from AutoShapes in xlsx via drawing XML.

    Returns {sheet_index: text} mapping.
    """
    import re
    import zipfile

    result = {}
    try:
        zf = zipfile.ZipFile(filepath)
    except zipfile.BadZipFile:
        return result

    # Build sheet_number -> drawing_file mapping from .rels
    sheet_to_drawing = {}
    for name in zf.namelist():
        if "worksheets/_rels/sheet" in name and name.endswith(".rels"):
            content = zf.read(name).decode("utf-8", errors="replace")
            for m in re.finditer(r'Target="([^"]*drawing[^"]*)"', content):
                # Extract sheet number from filename
                sheet_match = re.search(r"sheet(\d+)", name)
                if sheet_match:
                    sheet_num = int(sheet_match.group(1))
                    drawing_path = "xl/drawings/" + m.group(1).split("/")[-1]
                    sheet_to_drawing[sheet_num] = drawing_path

    # Extract text from each drawing
    for sheet_num, drawing_path in sheet_to_drawing.items():
        if drawing_path not in zf.namelist():
            continue
        content = zf.read(drawing_path).decode("utf-8", errors="replace")
        texts = re.findall(r"<a:t>([^<]+)</a:t>", content)
        if texts:
            # sheet number (1-based) to index (0-based)
            sheet_idx = sheet_num - 1
            if 0 <= sheet_idx < sheet_count:
                result[sheet_idx] = " ".join(texts)

    zf.close()
    return result


def _extract_pdf(filepath: Path) -> list[str]:
    import io
    from pdfminer.pdfpage import PDFPage
    from pdfminer.pdfinterp import PDFResourceManager, PDFPageInterpreter
    from pdfminer.converter import TextConverter
    from pdfminer.layout import LAParams

    pages = []
    rsrcmgr = PDFResourceManager()
    laparams = LAParams()
    with open(filepath, "rb") as f:
        for page in PDFPage.get_pages(f):
            output = io.StringIO()
            device = TextConverter(rsrcmgr, output, laparams=laparams)
            interpreter = PDFPageInterpreter(rsrcmgr, device)
            interpreter.process_page(page)
            pages.append(output.getvalue())
            device.close()
            output.close()
    return pages


def render_page_images(
    pdf_path: str | Path, output_dir: str | Path, doc_id: int
) -> list[Path]:
    """Render each page of a PDF as a PNG thumbnail (width ~600px).

    Saves to output_dir/<doc_id>/<page_num>.png (1-indexed).
    Returns a list of generated image paths.
    """
    import fitz  # PyMuPDF

    pdf_path = Path(pdf_path)
    output_dir = Path(output_dir) / str(doc_id)
    output_dir.mkdir(parents=True, exist_ok=True)

    # Suppress MuPDF warnings (e.g. structure tree errors from LibreOffice PDFs)
    fitz.TOOLS.mupdf_display_errors(False)
    fitz.TOOLS.mupdf_display_warnings(False)
    try:
        doc = fitz.open(pdf_path)
        paths = []
        for page_num in range(len(doc)):
            page = doc[page_num]
            # Scale to ~600px width
            zoom = 600.0 / page.rect.width if page.rect.width > 0 else 1.0
            mat = fitz.Matrix(zoom, zoom)
            pix = page.get_pixmap(matrix=mat)
            out_path = output_dir / f"{page_num + 1}.png"
            pix.save(str(out_path))
            paths.append(out_path)
        doc.close()
    finally:
        fitz.TOOLS.mupdf_display_errors(True)
        fitz.TOOLS.mupdf_display_warnings(True)
    return paths


def _chunk_text(text: str, chunk_size: int) -> list[str]:
    """Split text into chunks at paragraph boundaries."""
    paragraphs = text.split("\n")
    chunks = []
    current = []
    current_len = 0
    for para in paragraphs:
        if current_len + len(para) > chunk_size and current:
            chunks.append("\n".join(current))
            current = []
            current_len = 0
        current.append(para)
        current_len += len(para) + 1
    if current:
        chunks.append("\n".join(current))
    return chunks
